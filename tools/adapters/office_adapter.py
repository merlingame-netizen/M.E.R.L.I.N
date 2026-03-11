"""Office adapter — PowerPoint (python-pptx + COM), OneNote (COM), Teams (URI/process)."""

from __future__ import annotations

import subprocess
import sys
from pathlib import Path
from typing import Any

_HERE = Path(__file__).resolve().parent
_ROOT = _HERE.parent
for _p in (_ROOT, _HERE):
    _s = str(_p)
    if _s not in sys.path:
        sys.path.insert(0, _s)

from adapters.base_adapter import BaseAdapter  # noqa: E402

# ── Constants ────────────────────────────────────────────────────────────────

PYTHON = "C:/Users/PGNK2128/AppData/Local/Programs/Python/Python312/python.exe"

_PPTX_MISSING_MSG = (
    "python-pptx is not installed. "
    "Run: pip install python-pptx"
)
_ONENOTE_CLOSED_MSG = (
    "OneNote COM object not available. "
    "Please open Microsoft OneNote and try again."
)
_WIN32COM_MISSING_MSG = (
    "pywin32 is not installed. "
    "Run: pip install pywin32"
)


# ── Adapter ──────────────────────────────────────────────────────────────────


class OfficeAdapter(BaseAdapter):
    """
    Adapter for Microsoft Office / M365 desktop tools.

    PowerPoint: python-pptx (read/create) + win32com (PDF export).
    OneNote:    win32com — requires OneNote to be running.
    Teams:      process check + URI launch (no Graph API needed).
    """

    def __init__(self) -> None:
        super().__init__("office")

    # ── Action registry ──────────────────────────────────────────────────────

    def list_actions(self) -> dict[str, str]:
        return {
            # PowerPoint
            "ppt-open":       "Inspect a .pptx file: slide count, titles, notes (kwargs: path)",
            "ppt-info":       "Get metadata: author, title, slide count, word count (kwargs: path)",
            "ppt-export-pdf": "Export .pptx to PDF via PowerPoint COM (kwargs: path, out=None)",
            "ppt-create":     "Create a PPTX from title+bullets (kwargs: title, bullets=[], out='output.pptx')",
            # OneNote
            "onenote-notebooks": "List open notebooks",
            "onenote-sections":  "List sections in a notebook (kwargs: notebook=None for all)",
            "onenote-pages":     "List pages in a section (kwargs: section)",
            "onenote-read":      "Read page content as text (kwargs: page_id or title)",
            "onenote-create":    "Create a page with HTML content (kwargs: section, title, body_html)",
            # Teams
            "teams-status": "Check if Microsoft Teams is running",
            "teams-chat":   "Open Teams chat with a contact (kwargs: to — name or email)",
        }

    # ── Dispatch ─────────────────────────────────────────────────────────────

    def run(self, action: str, **kwargs: Any) -> dict:
        dispatch = {
            # PowerPoint
            "ppt-open":          self._ppt_open,
            "ppt-info":          self._ppt_info,
            "ppt-export-pdf":    self._ppt_export_pdf,
            "ppt-create":        self._ppt_create,
            # OneNote
            "onenote-notebooks": self._onenote_notebooks,
            "onenote-sections":  self._onenote_sections,
            "onenote-pages":     self._onenote_pages,
            "onenote-read":      self._onenote_read,
            "onenote-create":    self._onenote_create,
            # Teams
            "teams-status":      self._teams_status,
            "teams-chat":        self._teams_chat,
        }
        handler = dispatch.get(action)
        if handler is None:
            raise NotImplementedError
        return handler(**kwargs)

    # ── PowerPoint (python-pptx) ──────────────────────────────────────────────

    def _require_pptx(self):
        """Import and return pptx module, or raise ImportError with install hint."""
        try:
            import pptx  # noqa: PLC0415
            return pptx
        except ImportError:
            return None

    def _ppt_open(self, path: str = "", **_) -> dict:
        """Inspect a .pptx file: slide count, per-slide titles and notes."""
        if not path:
            return self.error("Missing required argument: path")
        pptx = self._require_pptx()
        if pptx is None:
            return self.error(_PPTX_MISSING_MSG)

        pptx_path = Path(path)
        if not pptx_path.exists():
            return self.error(f"File not found: {path}")

        self.log(f"Opening: {pptx_path.name}")
        try:
            prs = pptx.Presentation(str(pptx_path))
        except Exception as exc:
            return self.error(f"Could not open presentation: {exc}")

        slides = []
        for idx, slide in enumerate(prs.slides):
            title = _extract_title(slide)
            notes_text = ""
            if slide.has_notes_slide:
                tf = slide.notes_slide.notes_text_frame
                notes_text = tf.text.strip() if tf else ""
            slides.append({
                "index": idx + 1,
                "title": title,
                "notes_preview": notes_text[:200] if notes_text else "",
            })

        return self.ok({
            "path": str(pptx_path),
            "slide_count": len(slides),
            "slides": slides,
        })

    def _ppt_info(self, path: str = "", **_) -> dict:
        """Get file metadata: author, title, slide count, total word count."""
        if not path:
            return self.error("Missing required argument: path")
        pptx = self._require_pptx()
        if pptx is None:
            return self.error(_PPTX_MISSING_MSG)

        pptx_path = Path(path)
        if not pptx_path.exists():
            return self.error(f"File not found: {path}")

        self.log(f"Reading metadata: {pptx_path.name}")
        try:
            prs = pptx.Presentation(str(pptx_path))
        except Exception as exc:
            return self.error(f"Could not open presentation: {exc}")

        core = prs.core_properties
        word_count = _count_words(prs)

        return self.ok({
            "path": str(pptx_path),
            "title": core.title or "",
            "author": core.author or "",
            "last_modified_by": core.last_modified_by or "",
            "created": str(core.created) if core.created else "",
            "modified": str(core.modified) if core.modified else "",
            "slide_count": len(prs.slides),
            "word_count": word_count,
        })

    def _ppt_export_pdf(self, path: str = "", out: str | None = None, **_) -> dict:
        """Export a .pptx to PDF using PowerPoint COM (win32com)."""
        if not path:
            return self.error("Missing required argument: path")

        try:
            import win32com.client  # noqa: PLC0415
        except ImportError:
            return self.error(_WIN32COM_MISSING_MSG)

        pptx_path = Path(path).resolve()
        if not pptx_path.exists():
            return self.error(f"File not found: {path}")

        out_path = Path(out).resolve() if out else pptx_path.with_suffix(".pdf")

        self.log(f"Exporting to PDF via COM: {pptx_path.name} → {out_path.name}")
        app = None
        try:
            app = win32com.client.Dispatch("PowerPoint.Application")
            app.Visible = False
            prs = app.Presentations.Open(str(pptx_path), WithWindow=False)
            try:
                # SaveAs format 32 = ppSaveAsPDF
                prs.SaveAs(str(out_path), 32)
                self.log("PDF saved successfully.")
            finally:
                prs.Close()
        except Exception as exc:
            return self.error(f"PowerPoint COM export failed: {exc}")
        finally:
            if app is not None:
                try:
                    app.Quit()
                except Exception:
                    pass

        return self.ok({
            "input": str(pptx_path),
            "output": str(out_path),
            "format": "PDF",
        })

    def _ppt_create(self, title: str = "Untitled", bullets: list[str] | None = None,
                    out: str = "output.pptx", **_) -> dict:
        """Create a simple PPTX with a title slide and optional bullet slide."""
        pptx = self._require_pptx()
        if pptx is None:
            return self.error(_PPTX_MISSING_MSG)

        bullet_list = bullets or []
        out_path = Path(out)

        self.log(f"Creating presentation: {out_path.name}")
        try:
            from pptx.util import Inches, Pt  # noqa: PLC0415

            prs = pptx.Presentation()

            # Slide 1 — Title slide
            title_layout = prs.slide_layouts[0]
            slide1 = prs.slides.add_slide(title_layout)
            slide1.shapes.title.text = title
            if slide1.placeholders and len(slide1.placeholders) > 1:
                slide1.placeholders[1].text = ""

            # Slide 2 — Bullet slide (if content provided)
            if bullet_list:
                bullet_layout = prs.slide_layouts[1]
                slide2 = prs.slides.add_slide(bullet_layout)
                slide2.shapes.title.text = title
                tf = slide2.placeholders[1].text_frame
                tf.clear()
                for idx, bullet in enumerate(bullet_list):
                    if idx == 0:
                        tf.paragraphs[0].text = bullet
                    else:
                        p = tf.add_paragraph()
                        p.text = bullet
                        p.level = 0

            prs.save(str(out_path))
            self.log(f"Saved: {out_path}")
        except Exception as exc:
            return self.error(f"Failed to create presentation: {exc}")

        return self.ok({
            "path": str(out_path.resolve()),
            "slide_count": 2 if bullet_list else 1,
            "title": title,
            "bullets": bullet_list,
        })

    # ── OneNote (win32com) ────────────────────────────────────────────────────

    def _onenote_connect(self):
        """Return (onenote_app, None) or (None, error_dict)."""
        try:
            import win32com.client  # noqa: PLC0415
        except ImportError:
            return None, self.error(_WIN32COM_MISSING_MSG)
        try:
            app = win32com.client.Dispatch("OneNote.Application")
            return app, None
        except Exception as exc:
            return None, self.error(f"{_ONENOTE_CLOSED_MSG} Detail: {exc}")

    def _onenote_notebooks(self, **_) -> dict:
        """List open notebooks."""
        app, err = self._onenote_connect()
        if err is not None:
            return err

        self.log("Fetching notebooks from OneNote COM.")
        try:
            import xml.etree.ElementTree as ET  # noqa: PLC0415, N817

            xml_str = app.GetHierarchy("", 0, "")  # HierarchyScope.hsNotebooks = 0
            root = ET.fromstring(xml_str)
            ns = _onenote_ns(root)
            notebooks = []
            for nb in root.findall(f"{ns}Notebook"):
                notebooks.append({
                    "id": nb.get("ID", ""),
                    "name": nb.get("name", ""),
                    "path": nb.get("path", ""),
                    "is_open": nb.get("isCurrentlyViewed", "false") == "true",
                })
        except Exception as exc:
            return self.error(f"Failed to list notebooks: {exc}")

        return self.ok({"notebooks": notebooks, "count": len(notebooks)})

    def _onenote_sections(self, notebook: str | None = None, **_) -> dict:
        """List sections. notebook can be an ID or name; if None, all sections."""
        app, err = self._onenote_connect()
        if err is not None:
            return err

        self.log("Fetching sections from OneNote COM.")
        try:
            import xml.etree.ElementTree as ET  # noqa: PLC0415, N817

            # HierarchyScope.hsSections = 2
            if notebook:
                # Try to resolve notebook name → ID
                nb_id = _resolve_notebook_id(app, notebook)
                if nb_id is None:
                    return self.error(f"Notebook not found: {notebook}")
                xml_str = app.GetHierarchy(nb_id, 2, "")
            else:
                xml_str = app.GetHierarchy("", 2, "")

            root = ET.fromstring(xml_str)
            ns = _onenote_ns(root)
            sections = []
            for sec in root.iter(f"{ns}Section"):
                sections.append({
                    "id": sec.get("ID", ""),
                    "name": sec.get("name", ""),
                    "read_only": sec.get("readOnly", "false") == "true",
                })
        except Exception as exc:
            return self.error(f"Failed to list sections: {exc}")

        return self.ok({"sections": sections, "count": len(sections)})

    def _onenote_pages(self, section: str = "", **_) -> dict:
        """List pages in a section (ID or name required)."""
        if not section:
            return self.error("Missing required argument: section")

        app, err = self._onenote_connect()
        if err is not None:
            return err

        self.log(f"Fetching pages for section: {section}")
        try:
            import xml.etree.ElementTree as ET  # noqa: PLC0415, N817

            sec_id = _resolve_section_id(app, section)
            if sec_id is None:
                return self.error(f"Section not found: {section}")

            # HierarchyScope.hsPages = 3
            xml_str = app.GetHierarchy(sec_id, 3, "")
            root = ET.fromstring(xml_str)
            ns = _onenote_ns(root)
            pages = []
            for page in root.iter(f"{ns}Page"):
                pages.append({
                    "id": page.get("ID", ""),
                    "name": page.get("name", ""),
                    "date_time": page.get("dateTime", ""),
                    "is_currently_viewed": page.get("isCurrentlyViewed", "false") == "true",
                })
        except Exception as exc:
            return self.error(f"Failed to list pages: {exc}")

        return self.ok({"pages": pages, "count": len(pages)})

    def _onenote_read(self, page_id: str = "", title: str = "", **_) -> dict:
        """Read page content as plain text. Provide page_id or title."""
        if not page_id and not title:
            return self.error("Provide at least one of: page_id, title")

        app, err = self._onenote_connect()
        if err is not None:
            return err

        self.log(f"Reading OneNote page: {page_id or title}")
        try:
            import xml.etree.ElementTree as ET  # noqa: PLC0415, N817

            # Resolve title → ID if needed
            if not page_id:
                page_id = _resolve_page_id(app, title)
                if page_id is None:
                    return self.error(f"Page not found by title: {title}")

            # GetPageContent returns OneNote XML
            xml_str = app.GetPageContent(page_id, "")
            root = ET.fromstring(xml_str)
            # Extract all text nodes (crude but reliable for plain-text)
            text_parts = [t.strip() for t in root.itertext() if t.strip()]
            content = "\n".join(text_parts)
        except Exception as exc:
            return self.error(f"Failed to read page: {exc}")

        return self.ok({
            "page_id": page_id,
            "title": title,
            "content": content,
            "char_count": len(content),
        })

    def _onenote_create(self, section: str = "", title: str = "New Page",
                        body_html: str = "", **_) -> dict:
        """Create a page with HTML body in the specified section."""
        if not section:
            return self.error("Missing required argument: section")

        app, err = self._onenote_connect()
        if err is not None:
            return err

        self.log(f"Creating page '{title}' in section: {section}")
        try:
            sec_id = _resolve_section_id(app, section)
            if sec_id is None:
                return self.error(f"Section not found: {section}")

            # OneNote expects XHTML content wrapped in <one:Page>
            page_xml = (
                '<?xml version="1.0"?>'
                '<one:Page xmlns:one="http://schemas.microsoft.com/office/onenote/2013/onenote">'
                f'<one:Title><one:OE><one:T><![CDATA[{title}]]></one:T></one:OE></one:Title>'
                '<one:Outline><one:OEChildren><one:OE>'
                f'<one:T><![CDATA[{body_html}]]></one:T>'
                '</one:OE></one:OEChildren></one:Outline>'
                '</one:Page>'
            )
            new_page_id = ""
            app.CreateNewPage(sec_id, new_page_id, 0)  # 0 = npsDefault
            self.log("Page created successfully.")
        except Exception as exc:
            return self.error(f"Failed to create page: {exc}")

        return self.ok({
            "section": section,
            "title": title,
            "status": "created",
        })

    # ── Teams ─────────────────────────────────────────────────────────────────

    def _teams_status(self, **_) -> dict:
        """Check if Microsoft Teams is running via tasklist."""
        self.log("Checking Teams process via tasklist.")
        try:
            result = subprocess.run(
                ["tasklist", "/FI", "IMAGENAME eq ms-teams.exe"],
                capture_output=True,
                text=True,
                encoding="utf-8",
                errors="replace",
                timeout=10,
            )
            output1 = result.stdout

            # New Teams (msteams.exe)
            result2 = subprocess.run(
                ["tasklist", "/FI", "IMAGENAME eq msteams.exe"],
                capture_output=True,
                text=True,
                encoding="utf-8",
                errors="replace",
                timeout=10,
            )
            output2 = result2.stdout
        except subprocess.TimeoutExpired:
            return self.error("tasklist timed out.")
        except Exception as exc:
            return self.error(f"Failed to run tasklist: {exc}")

        running_classic = "ms-teams.exe" in output1.lower()
        running_new = "msteams.exe" in output2.lower()
        running = running_classic or running_new

        return self.ok({
            "running": running,
            "process": (
                "ms-teams.exe" if running_classic
                else "msteams.exe" if running_new
                else None
            ),
            "message": (
                "Microsoft Teams is running."
                if running
                else "Microsoft Teams is NOT running."
            ),
        })

    def _teams_chat(self, to: str = "", **_) -> dict:
        """Open a Teams chat with a contact using the ms-teams:// URI scheme."""
        if not to:
            return self.error("Missing required argument: to (name or email)")

        uri = f"msteams://teams.microsoft.com/l/chat/0/0?users={to}"
        self.log(f"Launching Teams chat URI for: {to}")
        try:
            # os.startfile / subprocess both work; subprocess is more controlled
            subprocess.Popen(
                ["cmd", "/c", "start", "", uri],
                creationflags=subprocess.CREATE_NO_WINDOW,
            )
        except Exception as exc:
            return self.error(f"Failed to launch Teams URI: {exc}")

        return self.ok({
            "to": to,
            "uri": uri,
            "message": f"Teams chat opened (or Teams window activated) for: {to}",
        })


# ── Private helpers ───────────────────────────────────────────────────────────


def _extract_title(slide) -> str:
    """Return slide title text or empty string."""
    if slide.shapes.title:
        return slide.shapes.title.text.strip()
    return ""


def _count_words(prs) -> int:
    """Count total words across all text frames in a presentation."""
    total = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                total += len(shape.text_frame.text.split())
    return total


def _onenote_ns(root) -> str:
    """Extract the XML namespace string from a OneNote root element."""
    tag = root.tag  # e.g. '{http://schemas.microsoft.com/office/onenote/2013/onenote}Notebooks'
    if tag.startswith("{"):
        return tag[:tag.index("}") + 1]
    return ""


def _resolve_notebook_id(app, name_or_id: str) -> str | None:
    """Find notebook ID by name (or pass through if already an ID)."""
    import xml.etree.ElementTree as ET  # noqa: N817
    try:
        xml_str = app.GetHierarchy("", 0, "")
        root = ET.fromstring(xml_str)
        ns = _onenote_ns(root)
        for nb in root.findall(f"{ns}Notebook"):
            if nb.get("ID") == name_or_id or nb.get("name") == name_or_id:
                return nb.get("ID")
    except Exception:
        pass
    return None


def _resolve_section_id(app, name_or_id: str) -> str | None:
    """Find section ID by name or pass through if already an ID."""
    import xml.etree.ElementTree as ET  # noqa: N817
    try:
        xml_str = app.GetHierarchy("", 2, "")
        root = ET.fromstring(xml_str)
        ns = _onenote_ns(root)
        for sec in root.iter(f"{ns}Section"):
            if sec.get("ID") == name_or_id or sec.get("name") == name_or_id:
                return sec.get("ID")
    except Exception:
        pass
    return None


def _resolve_page_id(app, title: str) -> str | None:
    """Find a page ID by title (searches all sections)."""
    import xml.etree.ElementTree as ET  # noqa: N817
    try:
        xml_str = app.GetHierarchy("", 3, "")
        root = ET.fromstring(xml_str)
        ns = _onenote_ns(root)
        for page in root.iter(f"{ns}Page"):
            if page.get("name") == title:
                return page.get("ID")
    except Exception:
        pass
    return None
